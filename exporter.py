"""
exporter.py - High-Stakes Consulting-Grade Word & PowerPoint export for TechM.
Designed for Billion Dollar Pitches.
"""
import os, re, sys, requests
from io import BytesIO
from datetime import datetime

if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")

from docx import Document
from docx.shared import Pt, RGBColor, Inches, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.enum.style import WD_STYLE_TYPE

from pptx import Presentation
from pptx.util import Inches as Pi, Pt as PPt
from pptx.dml.color import RGBColor as PR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Brand Colours (Tech Mahindra Consulting) ────────────────────────────────
TM_RED    = RGBColor(227, 24, 55)
TM_BLUE   = RGBColor(0, 32, 91)
TM_DBLUE  = RGBColor(0, 15, 45)
TM_GREY   = RGBColor(105, 105, 105)
TM_LGREY  = RGBColor(240, 242, 245)
TM_WHITE  = RGBColor(255, 255, 255)

PR_RED    = PR(227, 24, 55)
PR_BLUE   = PR(0, 32, 91)
PR_DBLUE  = PR(0, 15, 45)
PR_GREY   = PR(105, 105, 105)
PR_LGREY  = PR(240, 242, 245)
PR_WHITE  = PR(255, 255, 255)
PR_ACCENT = PR(200, 210, 230)

# ── Shared Helpers ─────────────────────────────────────────────────────────
def _sections(md):
    secs=[]; pat=re.compile(r"^#{1,2}\s+(.+)$",re.M); ms=list(pat.finditer(md))
    for i,m in enumerate(ms):
        body=md[m.end(): ms[i+1].start() if i+1<len(ms) else len(md)].strip()
        mer=re.search(r"```mermaid\s*([\s\S]+?)```",body,re.I)
        if not mer:
            mer_fallback = re.search(r"(flowchart\s+[TDBLRC]+|graph\s+[TDBLRC]+)[\s\S]+?(?:\n#{1,2}|\Z)", body, re.I)
            mermaid_code = mer_fallback.group(0).strip() if mer_fallback else None
        else:
            mermaid_code = mer.group(1).strip()
            
        # Sanitize common AI hallucinations in Mermaid
        if mermaid_code:
            mermaid_code = mermaid_code.replace("|>", "|")
            
        secs.append({
            "title": m.group(1).strip(),
            "content": re.sub(r"```mermaid[\s\S]+?```","",body,flags=re.I).strip(),
            "mermaid": mermaid_code
        })
    return secs

def _diagram(code):
    try:
        r = requests.post("https://kroki.io/mermaid/png", 
                          json={"diagram_source": code}, timeout=25)
        if r.status_code == 200 and len(r.content) > 800:
            return r.content
        else:
            print("Kroki error:", r.status_code, r.text[:100])
    except Exception as e:
        print("Diagram err:", e)
    return None

def _bullets(text):
    out = []
    for ln in text.split("\n"):
        ln = ln.strip()
        if not ln or ln.startswith("|") or re.match(r"^[-|:]+$", ln): continue
        ln = re.sub(r"^#{1,6}\s+", "", ln)
        ln = re.sub(r"^[-*•]\s+", "", ln)
        ln = re.sub(r"^\d+\.\s+", "", ln)
        ln = re.sub(r"\*\*(.+?)\*\*", r"\1", ln)
        ln = re.sub(r"\*(.+?)\*", r"\1", ln)
        ln = re.sub(r"`(.+?)`", r"\1", ln)
        if ln: out.append(ln)
    return out

def _parse_tables(text):
    tables = []; cur = []
    for ln in text.split("\n"):
        ln = ln.strip()
        if ln.startswith("|"):
            cur.append(ln)
        else:
            if cur: tables.append(cur); cur = []
    if cur: tables.append(cur)
    result = []
    for t in tables:
        rows = []
        for ln in t:
            if re.match(r"^\|[-| :]+\|$", ln): continue
            rows.append([c.strip() for c in ln.strip("|").split("|")])
        if rows: result.append(rows)
    return result

# ═══════════════════════════════════════════════════════════════════════════════
#  CONSULTING-GRADE WORD EXPORT
# ═══════════════════════════════════════════════════════════════════════════════
def _shd(cell, hex_color):
    tc = cell._tc; tcPr = tc.get_or_add_tcPr()
    s = OxmlElement("w:shd"); s.set(qn("w:fill"), hex_color)
    s.set(qn("w:color"), "auto"); s.set(qn("w:val"), "clear"); tcPr.append(s)

def _border_para(doc, color_hex="E31837", size_pt=2):
    p = OxmlElement("w:pBdr"); bb = OxmlElement("w:bottom")
    bb.set(qn("w:val"), "single"); bb.set(qn("w:sz"), str(int(size_pt*8)))
    bb.set(qn("w:space"), "1"); bb.set(qn("w:color"), color_hex)
    p.append(bb); return p

def _doc_heading(doc, text, level, color=TM_BLUE):
    h = doc.add_heading(text, level=level)
    if h.runs: 
        h.runs[0].font.color.rgb = color
        h.runs[0].font.name = 'Arial'
    if level == 1:
        pPr = h.paragraph_format._element.get_or_add_pPr()
        pPr.append(_border_para(doc))
    return h

def _doc_table(doc, rows):
    if not rows: return
    nc = max(len(r) for r in rows)
    t = doc.add_table(rows=len(rows), cols=nc)
    t.style = "Table Grid"
    t.autofit = True
    for i, row in enumerate(rows):
        cells = t.rows[i].cells
        for j in range(nc):
            txt = row[j] if j < len(row) else ""
            cells[j].text = txt
            para = cells[j].paragraphs[0]
            if not para.runs: para.add_run(txt)
            for run in para.runs:
                run.font.size = Pt(10)
                run.font.name = 'Arial'
                if i == 0: 
                    run.font.bold = True
                    run.font.color.rgb = TM_WHITE
            if i == 0: _shd(cells[j], "00205B")  # TM_BLUE
            elif i % 2 == 1: _shd(cells[j], "F0F2F5") # TM_LGREY
            
            # Vertical alignment
            tc = cells[j]._tc
            tcPr = tc.get_or_add_tcPr()
            vAlign = OxmlElement('w:vAlign')
            vAlign.set(qn('w:val'), 'center')
            tcPr.append(vAlign)

def export_to_word(outputs, client_name="Client", filename=None):
    os.makedirs("outputs", exist_ok=True)
    if not filename:
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"outputs/TechM_Consulting_{client_name.replace(' ','_')}_{ts}.docx"

    doc = Document()
    
    # Base styling
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Arial'
    font.size = Pt(11)
    font.color.rgb = RGBColor(50, 50, 50)
    
    # Layout and Margins
    for sec in doc.sections:
        sec.top_margin = Inches(1)
        sec.bottom_margin = Inches(1)
        sec.left_margin = Inches(1)
        sec.right_margin = Inches(1)
        
        # Header
        hdr = sec.header
        hp = hdr.paragraphs[0] if hdr.paragraphs else hdr.add_paragraph()
        hp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        r = hp.add_run("TECH MAHINDRA CONSULTING | HIGHLY CONFIDENTIAL")
        r.font.size = Pt(8)
        r.font.color.rgb = TM_GREY
        r.font.bold = True
        
        # Footer
        ftr = sec.footer
        fp = ftr.paragraphs[0] if ftr.paragraphs else ftr.add_paragraph()
        fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r2 = fp.add_run(f"Proprietary and Confidential  |  Prepared strictly for {client_name}  |  {datetime.now().strftime('%B %Y')}")
        r2.font.size = Pt(8)
        r2.font.color.rgb = TM_GREY

    # ── High-End Cover Page ──────────────────────────────────────────────────
    doc.add_paragraph("\n\n\n\n")
    
    logo_path = "techm_logo.jpg"
    if os.path.exists(logo_path):
        p_logo = doc.add_paragraph()
        p_logo.alignment = WD_ALIGN_PARAGRAPH.LEFT
        r_logo = p_logo.add_run()
        r_logo.add_picture(logo_path, width=Inches(3.5))
    else:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.LEFT
        r = p.add_run("TECH MAHINDRA")
        r.font.size = Pt(38)
        r.font.bold = True
        r.font.color.rgb = TM_RED
        r.font.name = 'Arial'

    p2 = doc.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.LEFT
    r2 = p2.add_run("STRATEGIC CONSULTING")
    r2.font.size = Pt(22)
    r2.font.bold = True
    r2.font.color.rgb = TM_BLUE
    r2.font.name = 'Arial'

    doc.add_paragraph("\n\n\n")
    p3 = doc.add_paragraph()
    p3.alignment = WD_ALIGN_PARAGRAPH.LEFT
    r3 = p3.add_run("Solution Proposal & Strategic Blueprint")
    r3.font.size = Pt(26)
    r3.font.color.rgb = TM_BLUE
    r3.font.bold = True

    p4 = doc.add_paragraph()
    p4.alignment = WD_ALIGN_PARAGRAPH.LEFT
    r4 = p4.add_run(f"Prepared exclusively for: {client_name}")
    r4.font.size = Pt(14)
    r4.font.color.rgb = TM_GREY

    doc.add_paragraph("\n\n\n\n\n\n")
    p5 = doc.add_paragraph()
    p5.alignment = WD_ALIGN_PARAGRAPH.LEFT
    r5 = p5.add_run("STRICTLY CONFIDENTIAL")
    r5.font.size = Pt(12)
    r5.font.bold = True
    r5.font.color.rgb = TM_RED

    p6 = doc.add_paragraph()
    p6.alignment = WD_ALIGN_PARAGRAPH.LEFT
    r6 = p6.add_run(f"Date: {datetime.now().strftime('%d %B %Y')}\nVersion: 1.0 (Final Draft)")
    r6.font.size = Pt(11)
    r6.font.color.rgb = TM_GREY
    
    doc.add_page_break()

    # ── Executive Summary & Body ──────────────────────────────────────────────
    proposal = outputs.get("proposal", "")
    secs = _sections(proposal)
    if not secs:
        doc.add_paragraph(proposal)
    else:
        for s in secs:
            _doc_heading(doc, s["title"], level=1)

            if s["mermaid"]:
                img = _diagram(s["mermaid"])
                if img:
                    p = doc.add_paragraph()
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    r = p.add_run()
                    r.add_picture(BytesIO(img), width=Inches(6.2))
                    
                    cap = doc.add_paragraph("Exhibit: High-Level Solution Architecture")
                    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    if cap.runs: 
                        cap.runs[0].italic = True
                        cap.runs[0].font.size = Pt(9)
                        cap.runs[0].font.color.rgb = TM_GREY

            if s["content"]:
                tables = _parse_tables(s["content"])
                clean = re.sub(r"\|.*\|", "", s["content"]).strip()
                for b in _bullets(clean):
                    p = doc.add_paragraph(style="List Bullet")
                    run = p.add_run(b)
                    run.font.size = Pt(11)
                for tbl in tables:
                    _doc_table(doc, tbl)
                    doc.add_paragraph()

    # ── Appendix ───────────────────────────────────────────────────────────────
    doc.add_page_break()
    _doc_heading(doc, "Appendix: Detailed Diagnostic Analysis", level=1, color=TM_GREY)
    for label, key in [("Requirements Analysis", "analysis"), ("Architecture Blueprint", "architecture"), ("Commercials & ROI", "pricing")]:
        content = outputs.get(key, "")
        if not content: continue
        _doc_heading(doc, label, level=2)
        tables = _parse_tables(content)
        clean = re.sub(r"\|.*\|", "", content).strip()
        for b in _bullets(clean)[:25]:
            p = doc.add_paragraph(style="List Bullet")
            p.add_run(b).font.size = Pt(10)
        for tbl in tables:
            _doc_table(doc, tbl)
            doc.add_paragraph()

    doc.save(filename)
    return filename

# ═══════════════════════════════════════════════════════════════════════════════
#  CONSULTING-GRADE POWERPOINT EXPORT
# ═══════════════════════════════════════════════════════════════════════════════
W = 13.33; H = 7.5  # 16:9 widescreen

def _bg(slide, rgb: PR):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = rgb

def _rect(slide, left, top, w, h, rgb: PR, line=False, shadow=False):
    s = slide.shapes.add_shape(1, Pi(left), Pi(top), Pi(w), Pi(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb
    if line:
        s.line.color.rgb = rgb
    else:
        s.line.fill.background()
    return s

def _tx(slide, text, left, top, w, h, size=16, bold=False,
        color=PR_WHITE, align=PP_ALIGN.LEFT, italic=False, name="Arial", valign=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Pi(left), Pi(top), Pi(w), Pi(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = PPt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = name
    return tb

def _footer(slide, slide_num=None, dark=False):
    color = PR_WHITE if dark else PR_GREY
    _tx(slide, "TECH MAHINDRA CONSULTING | HIGHLY CONFIDENTIAL", 
        0.5, H-0.4, 10, 0.3, size=8, color=color, align=PP_ALIGN.LEFT, bold=True)
    if slide_num:
        _tx(slide, str(slide_num), W-0.8, H-0.4, 0.4, 0.3, size=9, color=color, align=PP_ALIGN.RIGHT, bold=True)
    # Red accent line at bottom
    _rect(slide, 0, H-0.1, W, 0.1, PR_RED)

def _header_bar(slide, title, subtitle=""):
    # Premium Header: Dark blue with a red accent strip
    _rect(slide, 0, 0, W, 1.2, PR_DBLUE)
    _rect(slide, 0, 1.2, W, 0.05, PR_RED)
    
    _tx(slide, title.upper(), 0.5, 0.2, W-1, 0.6, size=24, bold=True, color=PR_WHITE)
    if subtitle:
        _tx(slide, subtitle, 0.5, 0.75, W-1, 0.4, size=12, color=PR_ACCENT, italic=True)

def _cover_slide(prs, client_name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, PR_DBLUE)
    
    # Left graphic accent
    _rect(slide, 0, 0, 0.8, H, PR_RED)
    _rect(slide, 0.8, 0, 0.2, H, PR_GREY)
    
    # Add TechM Logo
    logo_path = "techm_logo.jpg"
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Pi(1.8), Pi(1.2), height=Pi(0.8))
    else:
        _tx(slide, "TECH MAHINDRA", 1.8, 1.8, 10, 1.2, size=52, bold=True, color=PR_WHITE, name="Arial Black")
        
    # Main Titles
    _tx(slide, "STRATEGIC CONSULTING", 1.8, 3.0, 10, 0.8, size=24, bold=True, color=PR_RED)
    
    _rect(slide, 1.8, 3.8, 8, 0.05, PR_GREY)
    
    _tx(slide, "SOLUTION PROPOSAL & ROADMAP", 1.8, 4.2, 10, 0.6, size=20, color=PR_WHITE)
    _tx(slide, f"Prepared Exclusively For: {client_name}", 1.8, 4.8, 10, 0.6, size=16, color=PR_ACCENT)
    
    _tx(slide, f"Date: {datetime.now().strftime('%B %d, %Y')}", 1.8, 6.0, 5, 0.5, size=12, color=PR_GREY)
    _tx(slide, "STRICTLY CONFIDENTIAL", 1.8, 6.4, 5, 0.5, size=12, bold=True, color=PR_RED)
    
    _footer(slide, dark=True)

def _divider_slide(prs, title, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, PR_DBLUE)
    _rect(slide, 0, H/2 - 1, W, 2, PR_RED)
    
    _tx(slide, f"SECTION {num:02d}", 1.0, H/2 - 0.8, 11, 0.5, size=18, color=PR_WHITE, bold=True)
    _tx(slide, title.upper(), 1.0, H/2 - 0.2, 11, 1.0, size=36, bold=True, color=PR_WHITE)
    _footer(slide, dark=True)

def _exec_summary_slide(prs, outputs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, PR_LGREY)
    _header_bar(slide, "Executive Summary", "Strategic Highlights & Value Proposition")
    
    proposal = outputs.get("proposal", "")
    bullets = _bullets(proposal)[:4]
    
    # Left column: Metrics/KPIs
    _rect(slide, 0.5, 1.6, 3.5, 5.2, PR_WHITE)
    _rect(slide, 0.5, 1.6, 3.5, 0.1, PR_RED)
    
    metrics = [
        ("AI-DRIVEN", "Methodology"),
        ("100%", "Requirements Met"),
        ("SCALABLE", "Architecture"),
        ("OPTIMIZED", "Commercials")
    ]
    for i, (val, lbl) in enumerate(metrics):
        ly = 2.0 + i*1.2
        _tx(slide, val, 0.5, ly, 3.5, 0.5, size=24, bold=True, color=PR_DBLUE, align=PP_ALIGN.CENTER)
        _tx(slide, lbl, 0.5, ly+0.5, 3.5, 0.3, size=12, color=PR_GREY, align=PP_ALIGN.CENTER)
        if i < 3:
            _rect(slide, 1.5, ly+1.0, 1.5, 0.02, PR_LGREY)
            
    # Right column: Key insights
    _rect(slide, 4.3, 1.6, 8.5, 5.2, PR_WHITE)
    _rect(slide, 4.3, 1.6, 8.5, 0.1, PR_DBLUE)
    _tx(slide, "KEY PROPOSAL INSIGHTS", 4.8, 2.0, 7.5, 0.5, size=16, bold=True, color=PR_DBLUE)
    
    for i, b in enumerate(bullets):
        ly = 2.8 + i * 0.9
        _rect(slide, 4.8, ly+0.1, 0.15, 0.15, PR_RED)
        _tx(slide, b[:180], 5.1, ly, 7.2, 0.8, size=14, color=PR_DBLUE)
        
    _footer(slide, 2)

def _content_slide(prs, title, bullets, slide_num=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, PR_WHITE)
    _header_bar(slide, title)
    
    # Two column layout if many bullets
    if len(bullets) > 5:
        col1 = bullets[:5]
        col2 = bullets[5:10]
        for i, b in enumerate(col1):
            ly = 1.8 + i*1.0
            _rect(slide, 0.6, ly+0.12, 0.15, 0.15, PR_RED)
            _tx(slide, b[:200], 0.9, ly, 5.5, 0.9, size=14, color=PR_DBLUE)
        for i, b in enumerate(col2):
            ly = 1.8 + i*1.0
            _rect(slide, 7.0, ly+0.12, 0.15, 0.15, PR_RED)
            _tx(slide, b[:200], 7.3, ly, 5.5, 0.9, size=14, color=PR_DBLUE)
    else:
        for i, b in enumerate(bullets[:7]):
            ly = 1.8 + i*0.8
            _rect(slide, 0.8, ly+0.12, 0.15, 0.15, PR_RED)
            _tx(slide, b[:250], 1.2, ly, 11.0, 0.7, size=16, color=PR_DBLUE)
            
    _footer(slide, slide_num)

def _diagram_slide(prs, title, mermaid_code, slide_num=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, PR_WHITE)
    _header_bar(slide, title, "Target Operating Model & Architecture Blueprint")
    
    img = _diagram(mermaid_code)
    if img:
        # Center the image beautifully
        slide.shapes.add_picture(BytesIO(img), Pi(1.5), Pi(1.5), Pi(10.3), Pi(5.3))
    else:
        _rect(slide, 1.5, 1.5, 10.3, 5.3, PR_LGREY)
        _tx(slide, "ARCHITECTURE SCHEMATIC UNAVAILABLE", 1.5, 4.0, 10.3, 1, size=18, bold=True, color=PR_GREY, align=PP_ALIGN.CENTER)
    
    _footer(slide, slide_num)

def _table_slide(prs, title, rows, slide_num=None):
    if not rows: return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, PR_WHITE)
    _header_bar(slide, title, "Commercial & Operational Metrics")
    
    nc = max(len(r) for r in rows)
    nr = min(len(rows), 12)
    rh = min(0.45, 5.0/nr)
    
    tbl = slide.shapes.add_table(nr, nc, Pi(0.5), Pi(1.8), Pi(12.3), Pi(rh*nr)).table
    for i in range(nr):
        for j in range(nc):
            cell = tbl.cell(i, j)
            txt = rows[i][j] if j < len(rows[i]) else ""
            cell.text = txt
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = PPt(12 if nr <= 8 else 10)
            run.font.name = "Arial"
            
            if i == 0:
                run.font.bold = True
                run.font.color.rgb = PR_WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = PR_DBLUE
            else:
                run.font.color.rgb = PR_DBLUE
                cell.fill.solid()
                cell.fill.fore_color.rgb = PR_LGREY if i % 2 == 0 else PR_WHITE
                
    _footer(slide, slide_num)

def _thankyou_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, PR_DBLUE)
    
    _rect(slide, 0, H/2 - 1.5, W, 3, PR_RED)
    _tx(slide, "THANK YOU", 0, H/2 - 0.8, W, 1.0, size=64, bold=True, color=PR_WHITE, align=PP_ALIGN.CENTER)
    _tx(slide, "We look forward to a successful partnership.", 0, H/2 + 0.3, W, 0.5, size=20, color=PR_WHITE, align=PP_ALIGN.CENTER)
    
    _tx(slide, "TECH MAHINDRA STRATEGIC CONSULTING", 0, H-1.0, W, 0.5, size=14, bold=True, color=PR_ACCENT, align=PP_ALIGN.CENTER)
    _tx(slide, "confidential | proprietary | strictly not for distribution", 0, H-0.5, W, 0.5, size=10, color=PR_GREY, align=PP_ALIGN.CENTER)


def export_to_ppt(outputs, client_name="Client", filename=None):
    os.makedirs("outputs", exist_ok=True)
    if not filename:
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"outputs/TechM_Consulting_{client_name.replace(' ','_')}_{ts}.pptx"

    prs = Presentation()
    prs.slide_width = Pi(W)
    prs.slide_height = Pi(H)
    
    proposal = outputs.get("proposal", "")
    secs = _sections(proposal)

    _cover_slide(prs, client_name)
    _exec_summary_slide(prs, outputs)

    sn = 3
    for i, s in enumerate(secs):
        _divider_slide(prs, s["title"], i+1); sn+=1
        
        if s["mermaid"]:
            _diagram_slide(prs, "Architecture Blueprint", s["mermaid"], sn); sn+=1
            
        trows = _parse_tables(s["content"])
        if trows:
            for t in trows:
                _table_slide(prs, s["title"], t, sn); sn+=1
                
        clean = re.sub(r"\|.*\|", "", s["content"]).strip()
        buls = _bullets(clean)
        if buls:
            for chunk in [buls[i:i+10] for i in range(0, len(buls), 10)]:
                _content_slide(prs, s["title"], chunk, sn); sn+=1

    _thankyou_slide(prs)
    prs.save(filename)
    return filename